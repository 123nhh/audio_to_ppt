import os
import re
import time
import json
import shutil
import threading
from concurrent.futures import ThreadPoolExecutor, as_completed
from io import BytesIO
from PIL import Image, ImageFilter, ImageEnhance
from mutagen import File
from mutagen.flac import FLAC
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# --- OpenAI 库 ---
try:
    from openai import OpenAI
except ImportError:
    print("[错误] 缺少 openai 库，无法清洗歌词。请运行: pip install openai")
    OpenAI = None

# ==========================================
# 配置初始化
# ==========================================
CONFIG_FILE = "ai_config.json"
DEFAULT_KEY_PLACEHOLDER = "sk-xxxxxxxxxxxxxxxxxxxxxxxxxxxxxxxx"

def init_ai_configuration():
    default_config = {
        "enabled": True,
        "api_key": DEFAULT_KEY_PLACEHOLDER,
        "base_url": "https://api.openai.com/v1",
        "model": "gpt-3.5-turbo",
        "max_retries": 3,
        "max_workers": 4
    }

    if os.path.exists(CONFIG_FILE):
        try:
            with open(CONFIG_FILE, "r", encoding="utf-8") as f:
                saved_config = json.load(f)
                updated = False
                for k, v in default_config.items():
                    if k not in saved_config:
                        saved_config[k] = v
                        updated = True
                if updated:
                    with open(CONFIG_FILE, "w", encoding="utf-8") as fw:
                        json.dump(saved_config, fw, indent=4, ensure_ascii=False)
                return saved_config
        except: pass

    if not os.path.exists(CONFIG_FILE):
        try:
            with open(CONFIG_FILE, "w", encoding="utf-8") as f:
                json.dump(default_config, f, indent=4, ensure_ascii=False)
        except: pass
    
    return default_config

AI_CONFIG = init_ai_configuration()

# ==========================================
# 核心逻辑
# ==========================================

print_lock = threading.Lock()
def safe_print(msg):
    with print_lock:
        print(msg)

def call_ai_to_clean_lyrics(raw_text, log_tag):
    if not AI_CONFIG["enabled"] or not OpenAI: return raw_text
    if len(raw_text) < 10: return raw_text

    client = OpenAI(api_key=AI_CONFIG["api_key"], base_url=AI_CONFIG["base_url"])
    
    system_prompt = "你是一个专业的歌词整理助手。"
    user_prompt = (
        "请严格执行以下操作：\n"
        "1. 如果歌词包含'纯音乐'、'Instrumental'或没有实际歌词内容，请仅回复: [PURE_MUSIC]\n"
        "2. 删除头部元数据（作词、作曲、编曲等）。\n"
        "3. 删除尾部宣传信息（统筹、出品、邮箱等）。\n"
        "4. [重要] 如果是外语歌且包含翻译（如中文翻译），请删除所有翻译内容，只保留原文歌词。\n"
        "5. 必须保留原有的换行格式，不要随意合并行。\n"
        "6. [关键] 如果单行歌词过长（超过18个字符，包括标点符号在内），请根据语义停顿在中间插入符号 '^' 以便后续强制换行（例如：'长句的前半部分^长句的后半部分'）。注意仅插入符号，不要直接回车。\n"
        "7. [关键] 如果单行歌词内出现点号（逗号、句号等）将该符号换为 '^' ，若为该*行*歌词最后一个标点则删去该标点符号，注意仅插入或删去符号，不要直接回车。\n"
        "8. [关键] 如果单行歌词内出现括号、双引号等则在后面添加符号 '^' ，若为该*行*歌词最后一个标点则不做改动，注意仅插入符号，不要直接回车。（例如：'它匍匐（将人们恶言吐露）^痛哭（斥责神明的残酷）'）\n"
        "9. 不要输出任何解释，只输出结果。\n\n"
        "待处理文本：\n" + raw_text
    )

    for attempt in range(AI_CONFIG["max_retries"]):
        try:
            response = client.chat.completions.create(
                model=AI_CONFIG["model"],
                messages=[{"role": "system", "content": system_prompt}, {"role": "user", "content": user_prompt}],
                temperature=0.1, timeout=20
            )
            content = response.choices[0].message.content.strip()
            safe_print(f"[{log_tag}] [AI] 清洗成功 (已智能分行)")
            return content
        except Exception as e:
            safe_print(f"[{log_tag}] [警告] AI调用失败 (尝试 {attempt+1}): {e}")
            time.sleep(1)
    return raw_text

class AudioToPPT:
    def __init__(self, audio_path, output_ppt_path):
        self.audio_path = os.path.abspath(audio_path)
        self.output_ppt_path = os.path.abspath(output_ppt_path)
        self.file_name = os.path.basename(audio_path)
        self.is_pure_music = False
        
        self.mem_bg = None
        self.mem_cover = None
        self.mem_mask_top = None
        self.mem_mask_bottom = None
        
        self.metadata = {'title': '未知标题', 'artist': '未知歌手', 'lyrics': [], 'cover_data': None}
        
        # --- 📐 布局参数 ---
        self.SLIDE_W_INCH = 13.333
        self.SLIDE_H_INCH = 7.5
        
        self.SLIDE_W = Inches(self.SLIDE_W_INCH)
        self.SLIDE_H = Inches(self.SLIDE_H_INCH)
        self.CENTER_Y = self.SLIDE_H / 2
        
        self.MASK_H_INCH = 2.2
        self.LINE_SPACING = Inches(1.35)
        
        # --- [歌词页] 核心布局计算 ---
        ratio_left = 0.4
        zone_left_width = self.SLIDE_W_INCH * ratio_left
        zone_right_width = self.SLIDE_W_INCH * (1 - ratio_left)
        
        # [歌词页] 左侧封面布局
        self.LYRIC_COVER_SIZE_VAL = 3.5
        self.LYRIC_COVER_SIZE = Inches(self.LYRIC_COVER_SIZE_VAL)
        cover_margin = (zone_left_width - self.LYRIC_COVER_SIZE_VAL) / 2
        self.LYRIC_COVER_LEFT = Inches(cover_margin)
        
        # [歌词页] 右侧歌词布局
        self.TEXTBOX_W = Inches(7.8) 
        text_margin_in_zone = (zone_right_width - 7.8) / 2
        self.TEXTBOX_X = Inches(zone_left_width + text_margin_in_zone)
        self.TEXTBOX_H = Inches(2.2) 
        
        # 字体样式
        self.STYLE_ACTIVE = {'size': 40, 'bold': True, 'color': (255, 255, 255)} 
        self.STYLE_NORMAL = {'size': 24, 'bold': False, 'color': (160, 160, 160)} 

    def _log(self, msg):
        tag = self.metadata.get('title', self.file_name)
        safe_print(f"[{tag}] {msg}")

    def parse_lyrics_lines(self, text_content):
        cleaned_lines = []
        if not text_content: return cleaned_lines
        lines = text_content.split('\n')
        pattern = re.compile(r'\[\d{1,3}:\d{2}(?:\.\d{1,3})?\]')
        for line in lines:
            line_content = re.sub(pattern, '', line).strip()
            if line_content: cleaned_lines.append(line_content)
        return cleaned_lines

    def extract_metadata(self):
        try:
            audio = File(self.audio_path)
            tags = audio.tags
            if tags:
                raw_title = str(tags.get('TITLE', tags.get('TIT2', ['未知标题']))[0])
                self.metadata['title'] = raw_title.replace("《", "").replace("》", "").strip()
                self.metadata['artist'] = str(tags.get('ARTIST', tags.get('TPE1', ['未知歌手']))[0])
                
                raw_lyrics_text = ""
                if isinstance(audio, FLAC):
                    raw_lyrics_text = tags.get('lyrics', tags.get('unsyncedlyrics', ['']))[0]
                elif tags and hasattr(tags, 'getall'): 
                     uslt = tags.getall('USLT')
                     if uslt: raw_lyrics_text = uslt[0].text

                if raw_lyrics_text:
                    if "纯音乐" in raw_lyrics_text or "Instrumental" in raw_lyrics_text:
                        self.is_pure_music = True
                        self._log("检测到纯音乐标记")
                    else:
                        ai_result = call_ai_to_clean_lyrics(raw_lyrics_text, self.metadata['title'])
                        if "[PURE_MUSIC]" in ai_result:
                            self.is_pure_music = True
                            self._log("AI 判定为纯音乐")
                        else:
                            final_lines = self.parse_lyrics_lines(ai_result)
                            self.metadata['lyrics'] = final_lines
                            if not final_lines: self.is_pure_music = True
                else:
                    self.is_pure_music = True
                
                if isinstance(audio, FLAC) and audio.pictures:
                    self.metadata['cover_data'] = audio.pictures[0].data
                elif hasattr(audio, 'tags') and 'APIC:' in audio.tags: 
                     for key in audio.tags.keys():
                         if key.startswith('APIC'):
                             self.metadata['cover_data'] = audio.tags[key].data
                             break
        except Exception as e:
            self._log(f"[警告] 元数据错误: {e}")

    def image_to_bytes(self, img_obj, format='JPEG', quality=95):
        bio = BytesIO()
        img_obj.save(bio, format=format, quality=quality)
        bio.seek(0)
        return bio

    def add_gradient_transparency(self, img, direction='bottom'):
        img = img.convert("RGBA")
        width, height = img.size
        gradient = Image.new('L', (1, height), color=255)
        fade_len = 120 
        for y in range(height):
            if direction == 'bottom':
                if y >= height - fade_len:
                    alpha = int(255 * (height - y) / fade_len)
                    gradient.putpixel((0, y), alpha)
            elif direction == 'top':
                if y < fade_len:
                    alpha = int(255 * y / fade_len)
                    gradient.putpixel((0, y), alpha)
        alpha_mask = gradient.resize((width, height))
        img.putalpha(alpha_mask)
        return img

    def prepare_images(self):
        if not self.metadata['cover_data']: return None
        try:
            img = Image.open(BytesIO(self.metadata['cover_data'])).convert("RGB")
            target_w, target_h = 1280, 720
            
            # 高清背景
            bg_final = img.resize((target_w, target_h), resample=Image.LANCZOS)
            bg_final = bg_final.filter(ImageFilter.GaussianBlur(radius=60)) 
            bg_final = ImageEnhance.Brightness(bg_final).enhance(0.3)
            
            self.mem_bg = self.image_to_bytes(bg_final, format='JPEG') 
            self.mem_cover = self.image_to_bytes(img, format='JPEG')
            
            # 遮罩计算
            mask_ratio = self.MASK_H_INCH / self.SLIDE_H_INCH
            mask_pixel_h = int(target_h * mask_ratio)
            
            mask_top_img = bg_final.crop((0, 0, target_w, mask_pixel_h))
            mask_top_img = self.add_gradient_transparency(mask_top_img, direction='bottom')
            self.mem_mask_top = self.image_to_bytes(mask_top_img, format='PNG') 
            
            mask_bottom_img = bg_final.crop((0, target_h - mask_pixel_h, target_w, target_h))
            mask_bottom_img = self.add_gradient_transparency(mask_bottom_img, direction='top')
            self.mem_mask_bottom = self.image_to_bytes(mask_bottom_img, format='PNG')
            
            return True
        except Exception as e:
            self._log(f"[跳过] 图片处理失败: {e}")
            return False

    def create_cover_slide(self, prs):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(self.mem_bg, 0, 0, width=self.SLIDE_W, height=self.SLIDE_H)
        
        ALBUM_COVER_SIZE_VAL = 4.8
        ALBUM_COVER_SIZE = Inches(ALBUM_COVER_SIZE_VAL)
        
        center_x = (self.SLIDE_W - ALBUM_COVER_SIZE) / 2
        center_y = Inches(0.6) 
        
        slide.shapes.add_picture(self.mem_cover, center_x, center_y, width=ALBUM_COVER_SIZE, height=ALBUM_COVER_SIZE)
        
        text_top = center_y + ALBUM_COVER_SIZE + Inches(0.5)
        tx_intro = slide.shapes.add_textbox(0, text_top, self.SLIDE_W, Inches(1.8))
        tx_intro.text_frame.word_wrap = True
        
        p_title = tx_intro.text_frame.add_paragraph()
        p_title.text = self.metadata['title'] 
        p_title.font.size = Pt(36) 
        p_title.font.bold = True
        p_title.font.color.rgb = RGBColor(255, 255, 255)
        p_title.alignment = PP_ALIGN.CENTER
        
        p_artist = tx_intro.text_frame.add_paragraph()
        p_artist.text = self.metadata['artist']
        p_artist.font.size = Pt(20)
        p_artist.font.color.rgb = RGBColor(200, 200, 200)
        p_artist.alignment = PP_ALIGN.CENTER
        
        return slide

    def generate_ppt(self):
        if os.path.exists(self.output_ppt_path):
            try: os.remove(self.output_ppt_path)
            except: pass

        prs = Presentation()
        prs.slide_width = self.SLIDE_W
        prs.slide_height = self.SLIDE_H

        if not self.prepare_images():
            self._log("无法生成图片资源，跳过")
            return False

        self.create_cover_slide(prs)

        if self.is_pure_music or not self.metadata['lyrics']:
            self._log(f"纯音乐模式：仅生成封面")
            try: prs.save(self.output_ppt_path)
            except Exception as e: self._log(f"保存失败: {e}")
            return True

        lyrics = self.metadata['lyrics']
        
        lyric_cover_top = (self.SLIDE_H - self.LYRIC_COVER_SIZE) / 2 - Inches(1.0)
        lyric_text_top = lyric_cover_top + self.LYRIC_COVER_SIZE + Inches(0.2)

        for current_idx in range(len(lyrics)):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            
            # 1. 绘制背景 (最底层)
            slide.shapes.add_picture(self.mem_bg, 0, 0, width=self.SLIDE_W, height=self.SLIDE_H)
            
            # 2. [关键顺序调整] 绘制歌词 (中间层 - 下)
            # 歌词先画，这样它会被后面的遮罩盖住，但会被最后的封面压住（如果重叠的话）
            for target_idx in range(len(lyrics)):
                raw_line_text = lyrics[target_idx]
                line_text = raw_line_text.replace('^', '\n').replace(' ^ ', '\n')
                
                offset = target_idx - current_idx
                
                pos_y = self.CENTER_Y + (offset * self.LINE_SPACING) - (self.TEXTBOX_H / 2)
                
                tb = slide.shapes.add_textbox(self.TEXTBOX_X, pos_y, self.TEXTBOX_W, self.TEXTBOX_H)
                tf = tb.text_frame
                tf.word_wrap = True 
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE 
                
                p = tf.add_paragraph()
                p.text = line_text
                p.alignment = PP_ALIGN.CENTER 
                
                if offset == 0:
                    text_len = len(line_text)
                    if text_len > 30: final_size = 28
                    elif text_len > 18: final_size = 32
                    else: final_size = self.STYLE_ACTIVE['size']
                    
                    p.font.size = Pt(final_size)
                    p.font.bold = self.STYLE_ACTIVE['bold']
                    p.font.color.rgb = RGBColor(*self.STYLE_ACTIVE['color'])
                else:
                    text_len = len(line_text)
                    if text_len > 20: final_size = 20
                    else: final_size = self.STYLE_NORMAL['size']
                        
                    p.font.size = Pt(final_size)
                    p.font.bold = self.STYLE_NORMAL['bold']
                    p.font.color.rgb = RGBColor(*self.STYLE_NORMAL['color'])

            # 3. [关键顺序调整] 绘制遮罩 (中间层 - 上)
            # 遮罩要盖在歌词上面，所以放在歌词后面画
            try:
                slide.shapes.add_picture(self.mem_mask_top, 0, 0, width=self.SLIDE_W, height=Inches(self.MASK_H_INCH))
                slide.shapes.add_picture(self.mem_mask_bottom, 0, self.SLIDE_H - Inches(self.MASK_H_INCH), width=self.SLIDE_W, height=Inches(self.MASK_H_INCH))
            except: pass

            # 4. [关键顺序调整] 绘制左侧封面和信息 (最顶层)
            # 最后画封面，确保它在所有图层（包括遮罩）的最上面，不会被遮挡
            slide.shapes.add_picture(self.mem_cover, self.LYRIC_COVER_LEFT, lyric_cover_top, width=self.LYRIC_COVER_SIZE, height=self.LYRIC_COVER_SIZE)
            
            info_box = slide.shapes.add_textbox(self.LYRIC_COVER_LEFT, lyric_text_top, self.LYRIC_COVER_SIZE, Inches(2.0))
            info_box.text_frame.word_wrap = True
            
            p_song = info_box.text_frame.add_paragraph()
            p_song.text = self.metadata['title']
            p_song.font.size = Pt(20)
            p_song.font.bold = True
            p_song.font.color.rgb = RGBColor(255, 255, 255)
            p_song.alignment = PP_ALIGN.CENTER
            
            p_art = info_box.text_frame.add_paragraph()
            p_art.text = self.metadata['artist']
            p_art.font.size = Pt(14)
            p_art.font.color.rgb = RGBColor(180, 180, 180)
            p_art.alignment = PP_ALIGN.CENTER

        self.create_cover_slide(prs)

        try:
            prs.save(self.output_ppt_path)
            return True
        except Exception as e:
            self._log(f"保存失败: {e}")
            return False

def process_single_audio(file_path, output_dir):
    start_time = time.time()
    result = {"success": False, "is_pure": False, "duration": 0}
    
    try:
        file_base_name = os.path.splitext(os.path.basename(file_path))[0]
        relative_output_path = os.path.join(output_dir, f"{file_base_name}.pptx")
        
        converter = AudioToPPT(file_path, os.path.abspath(relative_output_path))
        converter.extract_metadata()
        
        if converter.generate_ppt():
            safe_print(f"[{file_base_name}] [完成] PPT已生成")
            result["success"] = True
            result["is_pure"] = converter.is_pure_music
        else:
            result["success"] = False
            
    except Exception as e:
        safe_print(f"[{file_path}] [失败] {e}")
        result["success"] = False
    
    result["duration"] = time.time() - start_time
    return result

def batch_process():
    input_dir, output_dir = "music", "output"
    if not os.path.exists(input_dir):
        os.makedirs(input_dir)
        print(f"[目录] 已自动创建 '{input_dir}' 文件夹。")
    if not os.path.exists(output_dir): os.makedirs(output_dir)

    root_files = [f for f in os.listdir('.') if os.path.isfile(f)]
    moved = 0
    for f in root_files:
        if f.lower().endswith(('.flac', '.mp3', '.wav', '.m4a')):
            try:
                shutil.move(os.path.abspath(f), os.path.join(os.path.abspath(input_dir), f))
                moved += 1
            except: pass
    if moved > 0: print(f"[整理] 已整理 {moved} 个文件到 music 文件夹。\n")

    files_in_dir = [f for f in os.listdir(input_dir) if f.lower().endswith(('.flac', '.mp3', '.wav', '.m4a'))]
    if not files_in_dir:
        print("[错误] music 文件夹为空。")
        return

    max_workers = AI_CONFIG.get("max_workers", 4)
    print(f"[开始] 发现 {len(files_in_dir)} 个文件 (并发: {max_workers})\n")
    
    total_start_time = time.time()
    stats = {
        "total_success": 0, "total_fail": 0,
        "pure_count": 0, "pure_duration": 0,
        "lyric_count": 0, "lyric_duration": 0
    }

    with ThreadPoolExecutor(max_workers=max_workers) as executor:
        futures = {executor.submit(process_single_audio, os.path.join(input_dir, f), output_dir): f for f in files_in_dir}
        
        for future in as_completed(futures):
            res = future.result()
            if res["success"]:
                stats["total_success"] += 1
                if res["is_pure"]:
                    stats["pure_count"] += 1
                    stats["pure_duration"] += res["duration"]
                else:
                    stats["lyric_count"] += 1
                    stats["lyric_duration"] += res["duration"]
            else:
                stats["total_fail"] += 1
    
    total_duration = time.time() - total_start_time
    avg_pure = stats["pure_duration"] / stats["pure_count"] if stats["pure_count"] > 0 else 0
    avg_lyric = stats["lyric_duration"] / stats["lyric_count"] if stats["lyric_count"] > 0 else 0

    print(f"\n" + "="*40)
    print(f"          详 细 统 计 报 告")
    print(f"="*40)
    print(f"[时间] 总耗时       : {total_duration:.2f} 秒")
    print(f"[成功] 处理成功     : {stats['total_success']} 首")
    print(f"[失败] 处理失败     : {stats['total_fail']} 首")
    print(f"-"*40)
    print(f"[音乐] 纯音乐       : {stats['pure_count']} 首")
    print(f"[速度] 纯音乐速度   : {avg_pure:.2f} 秒/首")
    print(f"-"*40)
    print(f"[歌词] 带歌词音乐   : {stats['lyric_count']} 首")
    print(f"[速度] 带歌词速度   : {avg_lyric:.2f} 秒/首")
    print(f"="*40 + "\n")
    print(f"[输出] 目录: {os.path.abspath(output_dir)}")

if __name__ == "__main__":
    batch_process()
